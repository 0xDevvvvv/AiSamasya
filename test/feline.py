from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a presentation object
prs = Presentation()

# Define slide content
slides_content = [
    {
    "title": "Data Link Layer Design Issues",
    "contents": [
      "Provides a well-defined service interface to the network layer.",
      "Deals with transmission errors.",
      "Regulates data flow (flow control) to prevent slow receivers from being overwhelmed by fast senders.",
      "Encapsulates network layer packets into frames (frame header, payload, frame trailer)."
    ],
    "bg_color": RGBColor(255,228,181),
    "title_color": RGBColor(255,165,0),
    "content_color": RGBColor(105,105,105)
  },
  {
    "title": "Data Link Layer Services",
    "contents": [
      "Unacknowledged Connectionless: Independent frames sent without acknowledgments (e.g., Ethernet).",
      "Acknowledged Connectionless: Each frame individually acknowledged (e.g., 802.11).",
      "Connection-Oriented: Connection established before data transfer, numbered frames, guaranteed delivery (in order and exactly once)."
    ],
    "bg_color": RGBColor(240,230,140),
    "title_color": RGBColor(255,215,0),
    "content_color": RGBColor(105,105,105)
  },
  {
    "title": "Framing Methods",
    "contents": [
      "Byte Count: Header specifies frame size (vulnerable to errors).",
      "Flag Bytes with Byte Stuffing: Frames start/end with special bytes, escape byte inserted before 'accidental' flag bytes in data.",
      "Flag Bits with Bit Stuffing: Frames start/end with special bit pattern, 0 bit inserted after five consecutive 1s.",
      "Physical Layer Coding Violations: Uses redundancy in physical layer encoding to delimit frames."
    ],
    "bg_color": RGBColor(255,235,205),
    "title_color": RGBColor(255,140,0),
    "content_color": RGBColor(105,105,105)
  },
  {
    "title": "Error Control",
    "contents": [
      "Positive/Negative acknowledgements.",
      "Timers to detect lost frames.",
      "Sequence numbers to identify retransmissions."
    ],
    "bg_color": RGBColor(255,248,220),
    "title_color": RGBColor(255,127,80),
    "content_color": RGBColor(105,105,105)
  },
  {
    "title": "Flow Control",
    "contents": [
      "Feedback-based: Receiver informs sender about its capacity.",
      "Rate-based: Protocol limits sender's transmission rate."
    ],
    "bg_color": RGBColor(245,222,179),
    "title_color": RGBColor(238,118,0),
    "content_color": RGBColor(105,105,105)
  },
  {
    "title": "Flow Control Techniques",
    "contents": [
      "Stop-and-Wait: Sender sends one frame, waits for acknowledgment before sending the next.",
      "Sliding Window: Sender sends multiple frames before waiting for acknowledgments, improving efficiency."
    ],
    "bg_color": RGBColor(250,240,210),
    "title_color": RGBColor(255,127,36),
    "content_color": RGBColor(105,105,105)
  },
  {
    "title": "ARQ Techniques",
    "contents": [
      "Stop-and-Wait ARQ: One-bit sliding window.",
      "Go-Back-N ARQ: Sender window size N, receiver window size 1, cumulative acknowledgments.",
      "Selective Repeat ARQ: Sender and receiver window sizes are equal, independent acknowledgments, only lost frames retransmitted."
    ],
    "bg_color": RGBColor(255,250,205),
    "title_color": RGBColor(210,105,30),
    "content_color": RGBColor(105,105,105)
  },
  {
    "title": "Data Link Protocols: HDLC",
    "contents": [
      "Bit-oriented protocol.",
      "Three station types: Primary, Secondary, Combined.",
      "Three frame types: Information (I), Supervisory (S), Unnumbered (U).",
      "Normal Response Mode (NRM) and Asynchronous Balanced Mode (ABM)."
    ],
    "bg_color": RGBColor(211,211,211),
    "title_color": RGBColor(128,128,128),
    "content_color": RGBColor(105,105,105)
  },
  {
    "title": "Data Link Protocols: PPP",
    "contents": [
      "Point-to-Point Protocol for Internet.",
      "Provides framing, link control (LCP), and network control (NCP).",
      "Supports multiple protocols, error detection, authentication, and optional reliable transmission."
    ],
    "bg_color": RGBColor(220,220,220),
    "title_color": RGBColor(160,160,160),
    "content_color": RGBColor(105,105,105)
  },
  {
    "title": "MAC Sublayer: IEEE 802",
    "contents": [
      "Handles medium access control in LANs/MANs.",
      "IEEE 802.3 (Ethernet): CSMA/CD.",
      "IEEE 802.4 (Token Bus): Token passing.",
      "IEEE 802.5 (Token Ring): Token passing."
    ],
    "bg_color": RGBColor(245,245,245),
    "title_color": RGBColor(192,192,192),
    "content_color": RGBColor(105,105,105)
  },
  {
    "title": "Ethernet (IEEE 802.3)",
    "contents": [
      "Classic Ethernet (CSMA/CD): 3-10 Mbps.",
      "Switched Ethernet: 100, 1000, 10000 Mbps (Fast Ethernet, Gigabit Ethernet, 10 Gigabit Ethernet).",
      "Frame format: Preamble, Destination Address, Source Address, Type, Data, Pad, Checksum.",
      "Addressing: Unicast, Multicast, Broadcast."
    ],
    "bg_color": RGBColor(230,230,230),
    "title_color": RGBColor(128,128,0),
    "content_color": RGBColor(105,105,105)
  },
  {
    "title": "IEEE 802.4 (Token Bus)",
    "contents": [
      "Stations arranged in a logical ring.",
      "Token holder allowed to transmit.",
      "Frame format includes preamble, start delimiter, frame control, addresses, payload, checksum, end delimiter."
    ],
    "bg_color": RGBColor(205,205,205),
    "title_color": RGBColor(139,69,19),
    "content_color": RGBColor(105,105,105)
  },
  {
    "title": "IEEE 802.5 (Token Ring)",
    "contents": [
      "Stations connected in a ring topology.",
      "Token passed between stations.",
      "Token holder permitted to transmit.",
      "Frame format includes start frame delimiter, access control, frame control, addresses, data, checksum, end delimiter, frame status."
    ],
    "bg_color": RGBColor(222,184,135),
    "title_color": RGBColor(139,69,19),
    "content_color": RGBColor(105,105,105)
  },
  {
    "title": "Bridges",
    "contents": [
      "Connect multiple LANs.",
      "Examine data link layer addresses for routing.",
      "Types: Routing bridges, Transparent bridges, Spanning Tree bridges, Remote bridges."
    ],
    "bg_color": RGBColor(255,218,185),
    "title_color": RGBColor(139,0,0),
    "content_color": RGBColor(105,105,105)
  },
  {
    "title": "Switches",
    "contents": [
      "Connect individual computers.",
      "Types: Two-layer switches (Cut Through), Three-layer switches (Routers)."
    ],
    "bg_color": RGBColor(255,192,203),
    "title_color": RGBColor(139,0,139),
    "content_color": RGBColor(105,105,105)
  },
  {
    "title": "High-Speed LANs: Fast Ethernet & Gigabit Ethernet",
    "contents": [
      "Fast Ethernet (IEEE 802.3u): 100 Mbps, backward-compatible with Standard Ethernet.",
      "Gigabit Ethernet (IEEE 802.3ab): 1 Gbps, full-duplex and half-duplex modes, jumbo frames."
    ],
    "bg_color": RGBColor(176,224,230),
    "title_color": RGBColor(0,0,139),
    "content_color": RGBColor(105,105,105)
  },
  {
    "title": "Wireless LANs: 802.11",
    "contents": [
      "Infrastructure mode: Clients connect to access points.",
      "Ad hoc mode: Clients connect directly to each other.",
      "Physical layer: Infrared, FHSS, DSSS, OFDM, HR-DSSS.",
      "MAC sublayer: CSMA/CA (Collision Avoidance), DCF (Distributed Coordination Function), PCF (Point Coordination Function).",     
      "Frame structure: Frame Control, Duration, Addresses, Sequence, Data, Checksum."
    ],
    "bg_color": RGBColor(210,180,140),
    "title_color": RGBColor(128,0,0),
    "content_color": RGBColor(105,105,105)
  },
  {
    "title": "802.15 (Bluetooth)",
    "contents": [
      "Wireless personal area network (PAN).",
      "Piconet: Master and up to 7 active slaves.",
      "Scatternet: Interconnected piconets.",
      "Protocol stack: Physical radio layer, baseband layer, link manager, L2CAP, audio/control protocols, middleware layer, application layer."
    ],
    "bg_color": RGBColor(190,190,190),
    "title_color": RGBColor(80,80,80),
    "content_color": RGBColor(105,105,105)
  }
]

# Function to add a slide
def add_slide(prs, title, content, bg_color, title_color, content_color):
    for line in content:
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg_color

        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = title_color
        title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        content_shape = slide.placeholders[1]
        content_shape.text = line
        for paragraph in content_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(32)
            paragraph.font.color.rgb = content_color
            paragraph.alignment = PP_ALIGN.CENTER

# Add slides to the presentation
for slide_content in slides_content:
    add_slide(prs, slide_content["title"], slide_content["contents"], slide_content["bg_color"], slide_content["title_color"], slide_content["content_color"])

# Save the presentation
prs.save('Waterfall_Model_Presentation.pptx')
