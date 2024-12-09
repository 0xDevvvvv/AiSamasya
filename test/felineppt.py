from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()

# Define slide content
slides_content = [
    {
        "title": "Waterfall Model",
        "content": "The stages of the waterfall model directly reflect the fundamental software development activities.",
        "bg_color": RGBColor(255, 255, 204)  # Light Yellow
    },
    {
        "title": "1. Requirements Analysis and Definition",
        "content": "The system’s services, constraints, and goals are established by consultation with system users. They are then defined in detail and serve as a system specification.",
        "bg_color": RGBColor(204, 255, 255)  # Light Cyan
    },
    {
        "title": "2. System and Software Design",
        "content": "The systems design process allocates the requirements to either hardware or software systems. It establishes an overall system architecture. Software design involves identifying and describing the fundamental software system abstractions and their relationships.",
        "bg_color": RGBColor(255, 204, 255)  # Light Pink
    },
    {
        "title": "3. Implementation and Unit Testing",
        "content": "During this stage, the software design is realized as a set of programs or program units. Unit testing involves verifying that each unit meets its specification.",
        "bg_color": RGBColor(204, 255, 204)  # Light Green
    },
    {
        "title": "4. Integration and System Testing",
        "content": "The individual program units or programs are integrated and tested as a complete system to ensure that the software requirements have been met. After testing, the software system is delivered to the customer.",
        "bg_color": RGBColor(255, 204, 204)  # Light Red
    },
    {
        "title": "5. Operation and Maintenance",
        "content": "Normally, this is the longest life-cycle phase. The system is installed and put into practical use. Maintenance involves correcting errors that were not discovered in earlier stages of the life cycle, improving the implementation of system units, and enhancing the system’s services as new requirements are discovered.",
        "bg_color": RGBColor(204, 204, 255)  # Light Purple
    },
    {
        "title": "Appropriate Use Cases",
        "content": "1. Embedded systems where the software has to interface with hardware systems.\n2. Critical systems where there is a need for extensive safety and security analysis of the software specification and design.\n3. Large software systems that are part of broader engineering systems developed by several partner companies.",
        "bg_color": RGBColor(255, 255, 153)  # Light Gold
    }
]

# Function to add a slide
def add_slide(prs, title, content, bg_color):
    slide_layout = prs.slide_layouts[5]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color

    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    content_shape = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8.5), Inches(5))
    text_frame = content_shape.text_frame
    p = text_frame.add_paragraph()
    p.text = content
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0, 0, 0)

# Add slides to the presentation
for slide_content in slides_content:
    add_slide(prs, slide_content["title"], slide_content["content"], slide_content["bg_color"])

# Save the presentation
prs.save('Waterfall_Model_Presentation.pptx')