from fileinput import filename 
from flask import *
import os
import PyPDF2
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import json
app = Flask(__name__) 
prs = Presentation()

genai.configure(api_key="AIzaSyAusjXdqax_JcIQMd5QBs_bLJGcc47KOAo")
model = genai.GenerativeModel("gemini-1.5-flash")

# Define the folder to save PowerPoint presentations
out_folder = 'output_ppts'

# Ensure the output folder exists
if not os.path.exists(out_folder):
    os.makedirs(out_folder)

app.config['out_folder'] = out_folder


def pdftotext(pdf_filename, txt_filename):
    with open(pdf_filename, 'rb') as pdf_file:
        reader = PyPDF2.PdfReader(pdf_file)
        text = ""
        
        for page_num in range(len(reader.pages)):
            page = reader.pages[page_num]
            text += page.extract_text()
        # print(text)
        return text
        # Write the extracted text to a .txt file
        # with open(txt_filename, 'w', encoding='utf-8') as txt_file:
        #     txt_file.write(text)

@app.route('/') 
def main(): 
    
    return render_template("index.html") 

def add_slide(prs, title, content, bg_color, title_color, content_color):
    for line in content:
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = convertrgb(bg_color)

        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = convertrgb(title_color)
        title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        content_shape = slide.placeholders[1]
        content_shape.text = line
        for paragraph in content_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(32)
            paragraph.font.color.rgb = convertrgb(content_color)
            paragraph.alignment = PP_ALIGN.CENTER

def convertrgb(color_list):
    r=int(color_list[0])
    g=int(color_list[1])
    b=int(color_list[2])
    return RGBColor(r,g,b)

@app.route('/output_ppts/<filename>')
def download(filename):
    return send_from_directory(app.config['out_folder'], filename, as_attachment=True)


# Save the generated presentation
def save_presentation(file_name):
    custom_filename = file_name.split('.')[0] + ".pptx"
    custom_file_path = os.path.join(out_folder, custom_filename)
    prs.save(custom_file_path)  # Save the presentation to the output folder
    return custom_filename


@app.route('/success', methods=['POST'])
def success(): 
    if request.method == 'POST': 
        f = request.files['file'] 
        f.save(f.filename) 
        pptname = f.filename.split('.')[0] + ".pptx"
        text = ""
        text += pdftotext(f.filename, 'output.txt')
        # Prompting
        prompt = text + """Summarise but make sure to include all the important content so that its easier 
        for an adhd person to understand, to the following JSON schema , give light vibrant colors such that
         the text is visible even with it
        with the rgb in the format [r,g,b] with no quotations at both ends:

        slide_content = {'title': str, 'contents': list[str],'bg_color':rgb(),'title_color':rgb(),'content_color':rgb()}
        Return: list[slide_content]"""
        response = model.generate_content(prompt)

        # Convert to JSON
        temp = response.text
        clean_temp = temp[temp.find("["):temp.rfind("]") + 1].strip()
        slides_content =json.loads(clean_temp)

        # Add slides to the presentation
        for slide_content in slides_content:
            add_slide(prs, slide_content["title"], slide_content["contents"], slide_content["bg_color"], slide_content["title_color"], slide_content["content_color"])

# Save the presentation
        save_presentation(pptname)

        os.remove(f.filename)
        return render_template("Acknowledgement.html", name=pptname) 
if __name__ == '__main__': 
    app.run()